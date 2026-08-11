# -*- coding: utf-8 -*-
"""
Graph API 권한 신청서 템플릿 및 테스트 샘플 PPTX 생성기

산출물
  Day2/templates/Graph API 권한 신청서_템플릿_v1.0.pptx   (배포용 빈 양식: Full / 약식)
  Day2/samples/*.pptx                                      (판정 엔진 테스트/데모용 작성 예시)

원본 기준: Graph API 할당 프로세스 및 기준_260806_학습용.pptx
  - 슬라이드 5  : 요청 템플렛(Full)
  - 슬라이드 11 : 요청 템플렛(약식)
"""
import io
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TPL_DIR = os.path.join(BASE, "templates")
SMP_DIR = os.path.join(BASE, "samples")
for d in (TPL_DIR, SMP_DIR):
    os.makedirs(d, exist_ok=True)

FONT = "맑은 고딕"
NAVY = RGBColor(0x1E, 0x27, 0x61)
INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x6B, 0x70, 0x7B)
LABEL_BG = RGBColor(0xE8, 0xEC, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 표 열 너비(합계 9.40인치) — 라벨/값 3쌍 구조
COL_W = [Inches(1.55), Inches(2.35), Inches(1.30), Inches(1.35), Inches(1.35), Inches(1.50)]

# ── 템플릿 행 정의 ────────────────────────────────────────────────────────────
# (라벨, 값키, span)  span=3 이면 값 셀을 3칸 병합
FULL_ROWS = [
    [("과제명", "taskName", 3), ("요청자/부서", "requester", 1)],
    [("과제개요", "taskOverview", 3), ("개발자/소속", "developer", 1)],
    [("*① 요청 API명/활용목적", "api1_np", 3), ("① Entra APP ID", "api1_appid", 1)],
    [("①활용 데이터/URL", "api1_data", 3), ("①Delegated/Application", "api1_type", 1)],
    [("② API명/활용목적", "api2_np", 3), ("② Entra APP ID", "api2_appid", 1)],
    [("②활용 데이터/URL", "api2_data", 3), ("②Delegated/Application", "api2_type", 1)],
    [("영향 범위/활용 규모", "scopeVolume", 1), ("정보주체 인지/동의", "consent", 1),
     ("저장/2차 가공", "storage", 1)],
    [("비고", "remark", 3), ("권한 만료일", "expiry", 1)],
]

SHORT_ROWS = [
    [("과제명", "taskName", 3), ("요청자/부서", "requester", 1)],
    [("과제개요", "taskOverview", 3), ("개발자/소속", "developer", 1)],
    [("*① API명/활용목적", "api1_np", 3), ("①Delegated/Application", "api1_type", 1)],
    [("①활용 데이터/URL", "api1_data", 3), ("①활용계정/App", "api1_account", 1)],
    [("영향 범위/활용 규모", "scopeVolume", 1), ("정보주체 인지/동의", "consent", 1),
     ("저장/2차 가공", "storage", 1)],
    [("*비고", "remark", 3), ("권한 만료일", "expiry", 1)],
]


def set_cell(cell, text, *, bold=False, label=False, size=9):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(45720)
    tf.margin_right = Emu(45720)
    tf.margin_top = Emu(18288)
    tf.margin_bottom = Emu(18288)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold or label
    run.font.color.rgb = NAVY if label else INK
    cell.fill.solid()
    cell.fill.fore_color.rgb = LABEL_BG if label else WHITE


def add_title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.30), Inches(0.20), Inches(9.40), Inches(0.42))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if sub:
        box2 = slide.shapes.add_textbox(Inches(0.32), Inches(0.66), Inches(9.40), Inches(0.30))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED


def add_note(slide, lines, top):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(top), Inches(9.36), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(8.5)
        r.font.color.rgb = MUTED


def build_table(slide, rows_def, values, top=Inches(1.05)):
    n_rows = len(rows_def)
    shape = slide.shapes.add_table(n_rows, 6, Inches(0.30), top, Inches(9.40), Inches(0.36 * n_rows))
    tbl = shape.table
    for i, w in enumerate(COL_W):
        tbl.columns[i].width = w

    for r, row_def in enumerate(rows_def):
        col = 0
        for label, key, span in row_def:
            set_cell(tbl.cell(r, col), label, label=True)
            val = values.get(key, "")
            vcell = tbl.cell(r, col + 1)
            if span == 3:
                # 값 셀 3칸 병합 (col+1 ~ col+3)
                vcell.merge(tbl.cell(r, col + 3))
            set_cell(vcell, val)
            col += 1 + span
    return tbl


def make_deck(path, slides):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec in slides:
        s = prs.slides.add_slide(blank)
        add_title(s, spec["title"], spec.get("sub"))
        rows = FULL_ROWS if spec["kind"] == "full" else SHORT_ROWS
        build_table(s, rows, spec["values"])
        note_top = 1.05 + 0.36 * len(rows) + 0.25
        add_note(s, spec.get("notes", []), note_top)
    prs.save(path)
    print("  saved:", os.path.basename(path))


GUIDE_FULL = [
    "* ① 요청 API명/활용목적 : Microsoft Graph 공식 권한명 + ' : ' + 활용목적  (예: Mail.Read : 특정 공유사서함의 알람 메일 수신 처리)",
    "* ①Delegated/Application : 로그인 사용자 권한 범위 내 사용이면 Delegated, 사용자 없이 앱 단독 실행이면 Application",
    "* 특정 사서함/팀/사이트만 사용하는 경우 활용목적·활용데이터에 그 대상을 반드시 명시 (등급 하향 및 범위 제한 조치 근거)",
    "* 권한 만료일은 필수 기재 (미기재 시 접수 보류). 최대 1년 이내 권장",
    "* 하나의 과제에 복수의 API가 필요한 경우 ③④… 행을 추가하여 기입",
]
GUIDE_SHORT = [
    "* IT Request 접수 시 기입 항목 : 권한 사용 App, Graph API명, 필요 데이터 필드, 활용 목적, 비고",
    "* ①활용계정/App : 실제 데이터를 읽거나 쓰는 대상 계정·사서함·사이트·팀을 구체적으로 기재",
    "* 저장/2차 가공 : 저장하지 않으면 'N/A', 저장 시 저장위치·보관기간·2차 활용범위를 함께 기재",
    "* 하나의 과제에 복수의 API가 필요한 경우 행을 추가하여 기입",
]

EMPTY_FULL = {k: "" for k in
              ["taskName", "taskOverview", "requester", "developer", "api1_np", "api1_appid",
               "api1_data", "api1_type", "api2_np", "api2_appid", "api2_data", "api2_type",
               "scopeVolume", "consent", "storage", "remark", "expiry"]}
EMPTY_SHORT = {k: "" for k in
               ["taskName", "taskOverview", "requester", "developer", "api1_np", "api1_type",
                "api1_data", "api1_account", "scopeVolume", "consent", "storage", "remark", "expiry"]}

print("[1] 배포용 빈 템플릿")
make_deck(
    os.path.join(TPL_DIR, "Graph API 권한 신청서_템플릿_v1.0.pptx"),
    [
        {"kind": "full", "title": "Graph API 권한 신청서 (Full)",
         "sub": "DX혁신실 M365 운영 | 접수 후 등급 분류(G1/G2/G3) 및 보안 검토가 진행됩니다",
         "values": EMPTY_FULL, "notes": GUIDE_FULL},
        {"kind": "short", "title": "Graph API 권한 신청서 (약식)",
         "sub": "IT Request 본문 기입용 약식 양식",
         "values": EMPTY_SHORT, "notes": GUIDE_SHORT},
    ],
)

# ── 테스트/데모 샘플 ─────────────────────────────────────────────────────────
SAMPLES = [
    # TC-1 : Application 신청 + '특정 사서함' → Delegated 재판정 → G1
    dict(file="샘플_TC1_특정사서함_알람메일.pptx", kind="full",
         title="Graph API 권한 신청서 (Full) — 작성 예시",
         sub="TC-1 : Application 신청이나 특정 공유사서함 한정 → Delegated 재판정 기대",
         values=dict(
             taskName="OCR 포털 시스템 알람 메일링",
             taskOverview="OCR 포털에서 발생하는 처리 결과/오류 알람을 배치로 취합하여 담당자에게 메일 발송",
             requester="윤찬희 대리/DX TF", developer="김ㅇㅇ/메가존",
             api1_np="Mail.Read : 특정 사서함을 통한 시스템 알람 메일링 목표",
             api1_appid="OCR-Portal-PRD-DP",
             api1_data="없음 (M365 데이터가 아닌 별도 포털 데이터의 메일 발송 기능만 활용) / noreply-ocr 공유사서함",
             api1_type="Application",
             api2_np="", api2_appid="", api2_data="", api2_type="",
             scopeVolume="OCR포털을 사용하는 전체 사용자 / Daily 배치를 통한 개별 알림 메일",
             consent="전원 동의/인지", storage="N/A",
             remark="공유사서함 noreply-ocr 1개만 사용", expiry="~27.1.31"),
         notes=GUIDE_FULL),

    # TC-2 : 전사 메일 본문 수집 → 재판정 미발동 → G3
    dict(file="샘플_TC2_전사메일수집.pptx", kind="full",
         title="Graph API 권한 신청서 (Full) — 작성 예시",
         sub="TC-2 : 전사 범위 민감 정보 → G3 유지 기대",
         values=dict(
             taskName="전사 메일 기반 AI 지식 검색",
             taskOverview="전사 메일 본문을 수집·적재하여 사내 지식 검색 서비스 학습 데이터로 활용",
             requester="이승목 과장/DX기획", developer="박ㅇㅇ/메가존",
             api1_np="Mail.Read : 전사 메일 본문을 수집하여 지식 검색 인덱스 구축",
             api1_appid="KM-Search-PRD",
             api1_data="메일 제목/본문/첨부파일명/발신자 (전체 사용자 사서함)",
             api1_type="Application",
             api2_np="Sites.ReadWrite.All : SharePoint 문서 수집 및 정제 결과 업로드",
             api2_appid="KM-Search-PRD",
             api2_data="전사 SharePoint 사이트 문서 목록 및 본문", api2_type="Application",
             scopeVolume="전 임직원 약 4,000명 / 매일 증분 수집 배치",
             consent="미정", storage="Azure AI Search 인덱스에 본문 적재, 벡터 임베딩 2차 가공",
             remark="", expiry=""),
         notes=GUIDE_FULL),

    # TC-5/6 : 2차 분류 하향/상향
    dict(file="샘플_TC5_본문제외_미저장.pptx", kind="full",
         title="Graph API 권한 신청서 (Full) — 작성 예시",
         sub="TC-5 : 데이터 최소화 + 미저장 → G3에서 1등급 하향(G2) 기대",
         values=dict(
             taskName="협업지수 산출 (HR Peer)",
             taskOverview="메일/채팅/회의 교신 '수'만 집계하여 조직 협업지수 산출",
             requester="김호연 대리/DX기획", developer="정ㅇㅇ/사내",
             api1_np="Mail.ReadWrite : 교신 건수 집계 및 결과 메일 발송",
             api1_appid="HR-Peer-PRD",
             api1_data="메일 본문 Field 제외 로직 구현, 발신자/수신자/일시 메타데이터만 사용 (마스킹 적용)",
             api1_type="Application",
             api2_np="", api2_appid="", api2_data="", api2_type="",
             scopeVolume="HR 평가 대상 임직원 1,200명 / 월 1회 집계",
             consent="전원 동의/인지", storage="N/A (저장 없이 통계 집계값만 활용)",
             remark="원문 미보관", expiry="2027-06-30"),
         notes=GUIDE_FULL),

    dict(file="샘플_TC6_임원Daily발송.pptx", kind="short",
         title="Graph API 권한 신청서 (약식) — 작성 예시",
         sub="TC-6 : 영향범위/규모 상향 요인 → G2에서 G3 상향 기대",
         values=dict(
             taskName="M365 활동지수 임원 리포트",
             taskOverview="직책자의 M365 활동 지수를 소속 임원에게 일 단위로 메일 발송",
             requester="강경현 차장/시스템운영", developer="한ㅇㅇ/사내",
             api1_np="User.Read.All : 직책자 프로필 및 조직 정보 조회",
             api1_type="Application",
             api1_data="표시이름, 부서, 직급, 상위조직 (전사 사용자 목록)",
             api1_account="전사 계정",
             scopeVolume="직책자 전원의 활동 지수를 소속 임원에게 Daily 메일 발송",
             consent="미정", storage="집계 결과 SharePoint 리스트 적재",
             remark="", expiry="2027-03-31"),
         notes=GUIDE_SHORT),

    # TC-8/10 : 고위험 + Delete
    dict(file="샘플_TC8_디렉터리쓰기_퇴직자정리.pptx", kind="short",
         title="Graph API 권한 신청서 (약식) — 작성 예시",
         sub="TC-8/10 : 권한상승 위험 + 삭제 작업 → G3 + Critical 기대",
         values=dict(
             taskName="퇴직자 계정/라이선스 자동 정리",
             taskOverview="HR 퇴직 확정 데이터를 연계하여 M365 계정 비활성화 및 라이선스 회수",
             requester="금동진 대리/시스템운영", developer="이ㅇㅇ/메가존",
             api1_np="Directory.ReadWrite.All : 퇴직 예정자 계정 비활성화 및 라이선스 제거",
             api1_type="Application",
             api1_data="사용자 계정 속성, 라이선스 할당 정보 (전사)",
             api1_account="전사 계정",
             scopeVolume="월 평균 30~50명 / 일 1회 배치",
             consent="인사 규정에 따름", storage="처리 이력 로그만 보관 (1년)",
             remark="", expiry="무기한"),
         notes=GUIDE_SHORT),

    # TC-14 : 범위 상충
    dict(file="샘플_TC14_범위상충.pptx", kind="short",
         title="Graph API 권한 신청서 (약식) — 작성 예시",
         sub="TC-14 : '특정 팀' + '전사' 동시 서술 → 재판정 미발동 + 범위 상충 기대",
         values=dict(
             taskName="Teams 공지 자동화",
             taskOverview="시스템 장애 발생 시 Teams 채널에 자동 공지",
             requester="한국남 대리/시스템운영", developer="최ㅇㅇ/사내",
             api1_np="ChannelMessage.Send : 특정 팀의 공지 채널에 장애 알림 게시 (추후 전사 확대)",
             api1_type="Application",
             api1_data="채널 메시지 본문 / 전사 모든 팀으로 확대 예정",
             api1_account="IT운영 공지팀 → 전사",
             scopeVolume="1차 IT운영팀 30명, 2차 전 임직원 / 장애 발생 시 실시간",
             consent="N/A", storage="N/A",
             remark="", expiry="2027-12-31"),
         notes=GUIDE_SHORT),
]

print("[2] 테스트/데모 샘플")
for s in SAMPLES:
    make_deck(os.path.join(SMP_DIR, s["file"]),
              [{"kind": s["kind"], "title": s["title"], "sub": s["sub"],
                "values": s["values"], "notes": s["notes"]}])

print("\n완료.")
print("  templates:", TPL_DIR)
print("  samples  :", SMP_DIR)
