# -*- coding: utf-8 -*-
"""
M365 관리자 역할 등급 판정 — 산출물 생성기

  A) PPT   : docs/M365 관리자 권한 분류 및 등급 기준_v1.0.pptx
  B) Excel : docs/M365 관리자 역할 명세서_v1.0.xlsx
  C) JSON  : tools/admin_roles.json   (HTML Tab2 임베드용)

원본 데이터: admin_roles_data.py
"""
import io
import json
import os
import sys
from collections import Counter, OrderedDict

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from admin_roles_data import (AXIS_LABEL, COMBO_MAP, GRADE_DEF, SENS_REASON, build)

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(BASE, "docs")
TOOLS = os.path.join(BASE, "tools")
os.makedirs(DOCS, exist_ok=True)

ROWS = build()
FONT = "맑은 고딕"
NAVY = RGBColor(0x1E, 0x27, 0x61)
INK = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x6B, 0x70, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HDR_BG = RGBColor(0xE8, 0xEC, 0xF6)
GC = {"A1": RGBColor(0x1B, 0x7F, 0x3B), "A2": RGBColor(0xB2, 0x6A, 0x00), "A3": RGBColor(0xB3, 0x26, 0x1E)}
GB = {"A1": RGBColor(0xE7, 0xF5, 0xEC), "A2": RGBColor(0xFE, 0xF3, 0xE2), "A3": RGBColor(0xFD, 0xEC, 0xEA)}

DIST = Counter(r["grade"] for r in ROWS)
CATS = OrderedDict()
for r in ROWS:
    CATS.setdefault(r["category"], []).append(r)


# ══════════════════════════════════════════════════════════════════════════
# C) JSON
# ══════════════════════════════════════════════════════════════════════════
def build_json():
    payload = {
        "meta": {"version": "1.0", "total": len(ROWS),
                 "source": "Microsoft Entra 기본 제공 역할 / Microsoft 365 관리 센터 관리자 역할",
                 "dist": dict(DIST)},
        "gradeDef": {g: {k: v for k, v in d.items()} for g, d in GRADE_DEF.items()},
        "comboMap": [{"combo": c, "scope": s, "op": o, "sens": n, "grade": g, "desc": d}
                     for c, s, o, n, g, d in COMBO_MAP],
        "axisLabel": AXIS_LABEL,
        "sensReason": SENS_REASON,
        "roles": ROWS,
    }
    p = os.path.join(TOOLS, "admin_roles.json")
    with open(p, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=1)
    print("  JSON  :", os.path.basename(p), f"({len(ROWS)}건)")
    return payload


def inject_html(payload):
    """HTML Tab2의 ADMIN 데이터 블록을 최신 데이터로 교체"""
    html_path = os.path.join(BASE, "GraphAPI_등급판정시스템.html")
    START = "/* ADMIN_ROLES_DATA_START"
    END = "/* ADMIN_ROLES_DATA_END */"
    with open(html_path, "r", encoding="utf-8") as f:
        html = f.read()
    i, j = html.find(START), html.find(END)
    if i < 0 or j < 0:
        raise RuntimeError("HTML에서 ADMIN_ROLES_DATA 마커를 찾을 수 없습니다.")
    block = (START + " — tools/build_admin_role_docs.py 가 이 블록을 갱신합니다. 직접 수정하지 마세요. */\n"
             "const ADMIN = " + json.dumps(payload, ensure_ascii=False, separators=(",", ":")) + ";\n")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html[:i] + block + html[j:])
    kb = (len(block.encode("utf-8"))) / 1024
    print("  HTML  : GraphAPI_등급판정시스템.html 에 역할 데이터 주입", f"({kb:.0f}KB)")


# ══════════════════════════════════════════════════════════════════════════
# B) Excel
# ══════════════════════════════════════════════════════════════════════════
XF = dict(
    hdr=PatternFill("solid", fgColor="1E2761"),
    sub=PatternFill("solid", fgColor="E8ECF6"),
    a1=PatternFill("solid", fgColor="E7F5EC"),
    a2=PatternFill("solid", fgColor="FEF3E2"),
    a3=PatternFill("solid", fgColor="FDECEA"),
)
THIN = Side(style="thin", color="C8CEDC")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def style_header(ws, row, ncols):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = XF["hdr"]
        cell.font = Font(name=FONT, size=9.5, bold=True, color="FFFFFF")
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = BORDER


def build_excel():
    wb = Workbook()

    # ── Sheet 1 : 판정 기준 ────────────────────────────────────────────────
    ws = wb.active
    ws.title = "판정 기준"
    ws.sheet_view.showGridLines = False
    ws["A1"] = "M365 / Microsoft Entra 관리자 역할 등급 판정 기준"
    ws["A1"].font = Font(name=FONT, size=16, bold=True, color="1E2761")
    ws["A2"] = "3기준 × 2값 = 8조합 → A1 / A2 / A3 · 사내 Graph API 등급 체계(G1/G2/G3)와 동일 구조"
    ws["A2"].font = Font(name=FONT, size=10, color="6B707B")
    ws["A3"] = f"총 {len(ROWS)}개 역할 · A3 {DIST['A3']}건 · A2 {DIST['A2']}건 · A1 {DIST['A1']}건"
    ws["A3"].font = Font(name=FONT, size=10, color="6B707B")

    r = 5
    ws.cell(r, 1, "① 3가지 판정 기준").font = Font(name=FONT, size=12, bold=True, color="1E2761")
    r += 1
    hdr = ["기준", "값", "코드", "정의"]
    for i, h in enumerate(hdr, 1):
        ws.cell(r, i, h)
    style_header(ws, r, len(hdr))
    axis_rows = [
        ("① 영향 범위 (Scope)", "서비스·기능 한정", "S", "단일 제품 또는 단일 기능 영역에만 적용"),
        ("", "테넌트 전역", "T", "디렉터리 전반(사용자·그룹·역할·앱·정책·도메인) 또는 전 서비스 교차"),
        ("② 작업 유형 (Operation)", "읽기 전용", "R", "조회·모니터링·보고만 가능"),
        ("", "구성 변경", "W", "생성·수정·삭제·정책 변경 가능"),
        ("③ 민감 자산 접근 (Sensitive)", "없음", "N", "권한 상승 경로도 콘텐츠 접근 경로도 없음"),
        ("", "있음", "P", "(a) 권한 상승 : 역할 할당·자격증명 변경·앱 권한 동의·CA/인증 정책·도메인/페더레이션 변경\n"
                        "(b) 콘텐츠 접근 : 메일·파일·채팅·문서 원문 열람 또는 열람 위임 설정 가능"),
    ]
    for a in axis_rows:
        r += 1
        for i, v in enumerate(a, 1):
            c = ws.cell(r, i, v)
            c.font = Font(name=FONT, size=9.5, bold=(i == 1 or i == 3))
            c.alignment = Alignment(vertical="center", wrap_text=True)
            c.border = BORDER
        ws.row_dimensions[r].height = 34 if a[2] == "P" else 20

    r += 2
    ws.cell(r, 1, "② 8조합 → 3등급 매핑").font = Font(name=FONT, size=12, bold=True, color="1E2761")
    r += 1
    hdr = ["조합", "범위", "작업", "민감 접근", "등급", "설명"]
    for i, h in enumerate(hdr, 1):
        ws.cell(r, i, h)
    style_header(ws, r, len(hdr))
    for combo, s, o, n, g, d in COMBO_MAP:
        r += 1
        vals = [combo, AXIS_LABEL[s], AXIS_LABEL[o], AXIS_LABEL[n], g, d]
        for i, v in enumerate(vals, 1):
            c = ws.cell(r, i, v)
            c.font = Font(name=FONT, size=9.5, bold=(i in (1, 5)))
            c.alignment = Alignment(horizontal="center" if i in (1, 5) else "left", vertical="center")
            c.border = BORDER
            c.fill = XF[g.lower()]

    r += 1
    ws.cell(r, 1, "핵심 : 민감 자산 접근(P)이 있으면 범위·작업과 무관하게 A3. "
                  "Graph API 기준의 'Delegated면 무조건 G1'과 동일한 단일축 우선 구조.")
    ws.cell(r, 1).font = Font(name=FONT, size=9.5, bold=True, color="B3261E")

    r += 2
    ws.cell(r, 1, "③ 등급별 승인 절차").font = Font(name=FONT, size=12, bold=True, color="1E2761")
    r += 1
    hdr = ["등급", "위험도", "구분", "정의", "승인 절차", "보안 검토", "PIM / 할당 방식", "예상 소요"]
    for i, h in enumerate(hdr, 1):
        ws.cell(r, i, h)
    style_header(ws, r, len(hdr))
    for g in ["A3", "A2", "A1"]:
        d = GRADE_DEF[g]
        r += 1
        vals = [g, d["risk"], d["label"], d["desc"], d["process"], d["review"], d["pim"], d["eta"]]
        for i, v in enumerate(vals, 1):
            c = ws.cell(r, i, v)
            c.font = Font(name=FONT, size=9.5, bold=(i == 1))
            c.alignment = Alignment(horizontal="center" if i in (1, 2) else "left",
                                    vertical="center", wrap_text=True)
            c.border = BORDER
            c.fill = XF[g.lower()]
        ws.row_dimensions[r].height = 40

    for col, w in zip("ABCDEFGH", [26, 20, 10, 62, 34, 30, 30, 12]):
        ws.column_dimensions[col].width = w

    # ── Sheet 2 : 역할 명세 ───────────────────────────────────────────────
    ws2 = wb.create_sheet("역할 명세")
    ws2.sheet_view.showGridLines = False
    cols = [
        ("권한 (역할명)", 30), ("영문 역할명", 44), ("위험도", 10), ("등급", 7),
        ("핵심 리스크", 76), ("범주", 13),
        ("범위", 14), ("작업", 10), ("민감 접근", 10), ("조합", 8),
        ("민감 접근 사유", 38), ("MS 권한있는역할", 14),
        ("보안 검토", 30), ("PIM / 할당 방식", 28), ("공식 문서", 48),
    ]
    for i, (h, _) in enumerate(cols, 1):
        ws2.cell(1, i, h)
    style_header(ws2, 1, len(cols))
    ws2.row_dimensions[1].height = 30

    for n, rr in enumerate(ROWS, start=2):
        g = rr["grade"]
        vals = [
            rr["ko"], rr["en"], rr["risk"], g, rr["keyRisk"], rr["category"],
            AXIS_LABEL[rr["scope"]], AXIS_LABEL[rr["op"]], AXIS_LABEL[rr["sens"]], rr["combo"],
            rr["sensReason"] or "-", "O" if rr["privileged"] else "",
            GRADE_DEF[g]["review"], GRADE_DEF[g]["pim"], rr["doc"],
        ]
        for i, v in enumerate(vals, 1):
            c = ws2.cell(n, i, v)
            c.font = Font(name=FONT, size=9, bold=(i in (1, 3, 4)))
            c.alignment = Alignment(horizontal="center" if i in (3, 4, 8, 9, 10, 12) else "left",
                                    vertical="center", wrap_text=(i == 5))
            c.border = BORDER
            if i in (3, 4):
                c.fill = XF[g.lower()]
                c.font = Font(name=FONT, size=9, bold=True, color=str(GC[g]))
        ws2.cell(n, 15).font = Font(name=FONT, size=8, color="3A5BA0", underline="single")
        ws2.cell(n, 15).hyperlink = rr["doc"]
        ws2.row_dimensions[n].height = 30

    for i, (_, w) in enumerate(cols, 1):
        ws2.column_dimensions[get_column_letter(i)].width = w
    ws2.freeze_panes = "C2"
    tbl = Table(displayName="RoleSpec", ref=f"A1:{get_column_letter(len(cols))}{len(ROWS)+1}")
    tbl.tableStyleInfo = TableStyleInfo(name="TableStyleLight1", showRowStripes=False)
    ws2.add_table(tbl)

    # ── Sheet 3 : 범주별 요약 ────────────────────────────────────────────
    ws3 = wb.create_sheet("범주별 요약")
    ws3.sheet_view.showGridLines = False
    ws3["A1"] = "범주별 · 등급별 역할 수"
    ws3["A1"].font = Font(name=FONT, size=14, bold=True, color="1E2761")
    hdr = ["범주", "A3 🔴 높음", "A2 🟡 중간", "A1 🟢 낮음", "합계"]
    for i, h in enumerate(hdr, 1):
        ws3.cell(3, i, h)
    style_header(ws3, 3, len(hdr))
    r = 3
    for cat, lst in sorted(CATS.items(), key=lambda kv: -len(kv[1])):
        r += 1
        cc = Counter(x["grade"] for x in lst)
        for i, v in enumerate([cat, cc["A3"], cc["A2"], cc["A1"], len(lst)], 1):
            c = ws3.cell(r, i, v)
            c.font = Font(name=FONT, size=10, bold=(i == 1 or i == 5))
            c.alignment = Alignment(horizontal="left" if i == 1 else "center", vertical="center")
            c.border = BORDER
    r += 1
    for i, v in enumerate(["합계", DIST["A3"], DIST["A2"], DIST["A1"], len(ROWS)], 1):
        c = ws3.cell(r, i, v)
        c.font = Font(name=FONT, size=10, bold=True, color="1E2761")
        c.alignment = Alignment(horizontal="left" if i == 1 else "center")
        c.border = BORDER
        c.fill = XF["sub"]
    for col, w in zip("ABCDE", [22, 14, 14, 14, 10]):
        ws3.column_dimensions[col].width = w

    ws3.cell(r + 3, 1, "운영 원칙").font = Font(name=FONT, size=12, bold=True, color="1E2761")
    principles = [
        "전역 관리자는 조직당 2~4명 이내로 제한하고, 상시 할당 대신 PIM 적격(Eligible) 할당으로 운영한다.",
        "전역 관리자 계정 잠금에 대비해 권한 있는 인증 관리자를 최소 1명 유지한다.",
        "모든 관리자 계정에 MFA를 필수 적용한다.",
        "작업 수행에 필요한 최소 권한 역할을 부여한다 (예: 암호 재설정만 필요하면 암호 관리자 / 기술 지원팀 관리자).",
        "A3 역할은 PIM 필수 · 상시 할당 금지 · 승인자 지정 · 활성화 사유 기록을 적용한다.",
        "분기별로 관리자 역할 할당 현황을 전수 조사하고, 미사용 할당은 회수한다.",
        "관리자 계정은 일반 업무 계정과 분리하여 운영한다.",
    ]
    for i, p in enumerate(principles):
        c = ws3.cell(r + 4 + i, 1, "· " + p)
        c.font = Font(name=FONT, size=9.5)

    p = os.path.join(DOCS, "M365 관리자 역할 명세서_v1.0.xlsx")
    wb.save(p)
    print("  Excel :", os.path.basename(p), f"(3시트 / {len(ROWS)}행)")


# ══════════════════════════════════════════════════════════════════════════
# A) PPT
# ══════════════════════════════════════════════════════════════════════════
def txbox(slide, x, y, w, h, text, size=10, bold=False, color=INK, align=PP_ALIGN.LEFT, space=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def slide_frame(prs, title, lead=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(1, Inches(0.12), Inches(0.13), Inches(9.76), Inches(0.50))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if lead:
        txbox(s, 0.34, 0.78, 9.32, 0.42, lead, size=10.5, bold=True, color=NAVY)
        ln = s.shapes.add_shape(1, Inches(0.34), Inches(1.19), Inches(9.32), Emu(9525))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0xC8, 0xCE, 0xDC)
        ln.line.fill.background()
    return s


def section(s, x, y, w, text):
    sh = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(0.26))
    sh.fill.solid()
    sh.fill.fore_color.rgb = HDR_BG
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = NAVY


def table(s, x, y, w, rows, widths, sizes=8.5, head_bg=NAVY, row_h=0.24, grade_col=None):
    n, m = len(rows), len(rows[0])
    shape = s.shapes.add_table(n, m, Inches(x), Inches(y), Inches(w), Inches(row_h * n))
    t = shape.table
    for i, cw in enumerate(widths):
        t.columns[i].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(45720)
            tf.margin_top = tf.margin_bottom = Emu(9144)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (ri == 0 or ci in (grade_col or [])) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            run.font.size = Pt(sizes if ri else sizes)
            run.font.bold = (ri == 0) or (ci in (grade_col or []))
            run.font.color.rgb = WHITE if ri == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = head_bg if ri == 0 else WHITE
            if ri and grade_col and ci in grade_col:
                g = str(val).strip()[:2]
                if g in GB:
                    cell.fill.fore_color.rgb = GB[g]
                    run.font.color.rgb = GC[g]
    return t


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1. 배경 및 목적 ──────────────────────────────────────────────────
    s = slide_frame(prs, "M365 관리자 권한 분류 및 등급 기준 수립 배경",
                    "관리자 역할 요청에 대한 신속하고 일관된 대응 체계 수립")
    section(s, 0.29, 1.30, 1.9, "관리자 역할 개요")
    txbox(s, 0.34, 1.66, 9.32, 0.9, [
        " 관리자 역할 : M365·Entra 관리 센터에서 특정 관리 작업을 수행할 수 있도록 부여하는 사전 정의 권한 집합",
        " 현황 : Microsoft Entra 기본 제공 역할 " + str(len(ROWS)) + "종. 역할마다 영향 범위와 위험도가 크게 상이",
        " 문제 : Graph API와 달리 관리자 역할은 판정 기준이 없어 요청 건마다 담당자 재량으로 승인 여부를 판단",
    ], size=9.5, space=3)

    section(s, 0.29, 2.68, 2.6, "Microsoft 공식 보안 지침")
    table(s, 0.31, 3.00, 9.38, [
        ["권장 사항", "중요한 이유"],
        ["가능한 한 적은 수의 전역 관리자 유지",
         "전역 관리자는 조직 설정과 대부분의 데이터에 거의 무제한 접근이 가능. 계정 잠금 대비 권한 있는 인증 관리자 최소 1명 유지"],
        ["최소 허용 역할 할당",
         "작업 완료에 필요한 접근 권한만 부여. 암호 재설정만 필요하면 전역 관리자가 아닌 암호 관리자·기술 지원팀 관리자를 할당"],
        ["관리자 MFA 필수",
         "관리자는 사용자 데이터에 접근 가능. 암호가 유출되어도 두 번째 인증 없이는 로그인 불가하도록 통제"],
    ], [3.0, 6.38], sizes=8.5, row_h=0.42)

    section(s, 0.29, 4.95, 2.2, "추진 목표 및 산출물")
    table(s, 0.31, 5.27, 9.38, [
        ["항목", "산출물 / 체계"],
        ["관리자 역할 분류 등급 기준 수립", "① 3기준 × 2값 = 8조합 → A1/A2/A3 3단계 분류  ② 역할별 위험도·핵심 리스크 명세"],
        ["요청 검토 / 승인 프로세스 수립", "① 등급별 승인 절차 및 R&R 정의  ② 보안 검토 대상 기준 명확화"],
        ["할당 이력 및 사후 관리 체계", "① PIM 기반 적격 할당 운영  ② 분기별 전수 조사 및 미사용 할당 회수"],
    ], [3.0, 6.38], sizes=8.5, row_h=0.34)
    txbox(s, 0.33, 6.95, 9.3, 0.5, [
        "* 사내 Graph API 권한 등급 체계(G1/G2/G3)와 동일한 '3기준 × 2값 = 8조합 → 3등급' 구조로 설계하여 판정 방식의 일관성 확보",
        "* 출처 : Microsoft 365 관리 센터 관리자 역할 정보 / Microsoft Entra 기본 제공 역할 (learn.microsoft.com)",
    ], size=7.5, color=MUTED, space=2)

    # ── 2. 1차 분류 기준 ─────────────────────────────────────────────────
    s = slide_frame(prs, "관리자 역할 등급 분류 기준 (1/2) — 3가지 판정 기준",
                    "관리자 역할을 ①얼마나 넓게 ②무엇을 하며 ③민감 자산에 닿는가 3가지 기준으로 분류")
    section(s, 0.29, 1.28, 2.4, "1차 분류 — 3기준 × 2값")
    table(s, 0.31, 1.58, 9.38, [
        ["구분", "", "정의", "예시 역할"],
        ["① 영향 범위\n(Scope)", "S. 서비스·기능 한정",
         "단일 제품 또는 단일 기능 영역에만 적용", "Exchange 관리자, Teams 관리자, 라이선스 관리자, 청구 관리자"],
        ["", "T. 테넌트 전역",
         "디렉터리 전반(사용자·그룹·역할·앱·정책·도메인) 또는 전 서비스 교차", "전역 관리자, 사용자 관리자, 보안 관리자, 전역 읽기 권한자"],
        ["② 작업 유형\n(Operation)", "R. 읽기 전용",
         "조회·모니터링·보고만 가능. 설정 변경 불가", "전역 읽기 권한자, 보안 읽기 권한자, 보고서 읽기 권한자"],
        ["", "W. 구성 변경",
         "생성·수정·삭제·정책 변경 가능", "대부분의 '○○ 관리자' 역할"],
        ["③ 민감 자산 접근\n(Sensitive)", "N. 없음",
         "권한 상승 경로도 콘텐츠 접근 경로도 없음", "라이선스 관리자, Teams 관리자, 프린터 관리자"],
        ["", "P. 있음",
         "(a) 권한 상승 : 역할 할당 · 자격 증명/MFA 변경 · 앱 권한 동의 · 조건부 액세스/인증 정책 · 도메인/페더레이션 변경\n"
         "(b) 콘텐츠 접근 : 메일 · 파일 · 채팅 · 문서 원문 열람 또는 열람 위임 설정 가능",
         "전역 관리자, 사용자 관리자, Exchange 관리자, SharePoint 관리자, Intune 관리자, 준수 관리자"],
    ], [1.35, 1.85, 3.55, 2.63], sizes=8, row_h=0.46, grade_col=[])

    section(s, 0.29, 5.28, 2.9, "민감 자산 접근(P) 판정 사유 코드")
    codes = list(SENS_REASON.items())
    half = (len(codes) + 1) // 2
    reason_rows = [["코드", "사유", "코드", "사유"]]
    for i in range(half):
        lf = codes[i]
        rt = codes[i + half] if i + half < len(codes) else ("", "")
        reason_rows.append([lf[0], lf[1], rt[0], rt[1]])
    table(s, 0.31, 5.58, 9.38, reason_rows, [1.20, 3.49, 1.20, 3.49], sizes=7.5, row_h=0.22)
    txbox(s, 0.33, 6.92, 9.3, 0.2,
          "* 위 두 경로(권한 상승 / 콘텐츠 접근) 중 하나라도 해당하면 P로 판정한다.",
          size=7.5, color=MUTED)

    # ── 3. 8조합 매핑 ────────────────────────────────────────────────────
    s = slide_frame(prs, "관리자 역할 등급 분류 기준 (2/2) — 8조합 → 3등급",
                    "3기준의 8가지 조합(2×2×2)을 A1 / A2 / A3 3개 등급으로 매핑")
    section(s, 0.29, 1.30, 2.0, "8조합 → 등급 매핑")
    combo_rows = [["조합", "범위", "작업", "민감 접근", "등급", "설명"]]
    for combo, sc, op, sn, g, d in COMBO_MAP:
        combo_rows.append([combo, AXIS_LABEL[sc], AXIS_LABEL[op], AXIS_LABEL[sn], g, d])
    table(s, 0.31, 1.62, 9.38, combo_rows, [0.80, 1.55, 1.00, 1.00, 0.72, 4.31],
          sizes=8.5, row_h=0.245, grade_col=[0, 4])
    txbox(s, 0.33, 3.90, 9.3, 0.36,
          "핵심 : 민감 자산 접근(P)이 있으면 범위·작업과 무관하게 A3. "
          "Graph API 기준에서 'Delegated면 무조건 G1'이 등급을 결정했듯, 관리자 역할은 '민감 자산 접근'이 등급을 결정한다.",
          size=9, bold=True, color=GC["A3"])

    section(s, 0.29, 4.32, 2.2, "등급별 정의 및 승인 절차")
    grade_rows = [["등급 / 구분", "위험도", "정의", "승인 절차", "보안 검토", "PIM / 할당", "소요"]]
    for g in ["A1", "A2", "A3"]:
        d = GRADE_DEF[g]
        grade_rows.append([g + "\n" + d["label"], d["risk"], d["desc"],
                           d["process"], d["review"], d["pim"], d["eta"]])
    table(s, 0.31, 4.62, 9.38, grade_rows, [1.10, 0.72, 2.66, 1.55, 1.42, 1.30, 0.63],
          sizes=7.5, row_h=0.55, grade_col=[0])

    dist_txt = (f"판정 결과 분포 (총 {len(ROWS)}개 역할) : "
                f"A3 {DIST['A3']}건 ({DIST['A3']*100//len(ROWS)}%) · "
                f"A2 {DIST['A2']}건 ({DIST['A2']*100//len(ROWS)}%) · "
                f"A1 {DIST['A1']}건 ({DIST['A1']*100//len(ROWS)}%)")
    txbox(s, 0.33, 6.98, 9.3, 0.5, [
        dist_txt,
        "* A3 비중이 높은 것은 관리자 역할 다수가 구조적으로 권한 상승 또는 콘텐츠 접근 경로를 갖기 때문이며, "
        "이는 관리자 역할 요청을 Graph API보다 엄격히 통제해야 하는 근거가 된다.",
    ], size=8, color=MUTED, space=3)

    # ── 4. 신청 / 승인 프로세스 ──────────────────────────────────────────
    s = slide_frame(prs, "관리자 권한 신청 및 승인 프로세스 (案)",
                    "등급에 따라 승인 절차와 할당 방식을 차등 적용")
    section(s, 0.29, 1.30, 1.8, "진행 절차")
    steps = [("1. 권한 요청 접수", "요청자 / DX기획"), ("2. 역할 등급 판정", "DX운영"),
             ("3. 보안 검토", "정보보호"), ("4. IR 결재", "요청자"),
             ("5. 역할 할당 (PIM)", "M365 운영"), ("6. 사후 모니터링", "DX운영")]
    x = 0.33
    for name, who in steps:
        sh = s.shapes.add_shape(5, Inches(x), Inches(1.66), Inches(1.50), Inches(0.62))
        sh.fill.solid()
        sh.fill.fore_color.rgb = HDR_BG
        sh.line.color.rgb = RGBColor(0xC8, 0xCE, 0xDC)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run(); r1.text = name
        r1.font.name = FONT; r1.font.size = Pt(8.5); r1.font.bold = True; r1.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = who
        r2.font.name = FONT; r2.font.size = Pt(7); r2.font.color.rgb = MUTED
        x += 1.57

    section(s, 0.29, 2.50, 2.6, "등급별 차등 적용 사항")
    table(s, 0.31, 2.82, 9.38, [
        ["구분", "A1 🟢 낮음", "A2 🟡 중간", "A3 🔴 높음"],
        ["보안 검토", "생략 (사후 이력 관리)", "정보보호그룹 공유 · 의견 수렴", "정보보호그룹 검토 필수\n필요시 Committee 상정"],
        ["할당 방식", "상시 할당 가능", "PIM 적격 할당 권장", "PIM 적격 할당 필수\n상시(Active) 할당 금지"],
        ["활성화 조건", "-", "MFA", "MFA + 승인자 승인 + 사유 기록\n최대 활성화 시간 8시간"],
        ["할당 기간", "1년 (연장 심사)", "6개월 (연장 심사)", "3개월 (연장 심사)"],
        ["계정 분리", "일반 계정 가능", "일반 계정 가능", "관리자 전용 계정 필수"],
        ["모니터링", "분기 전수 조사", "분기 전수 조사", "분기 전수 조사 + 활성화 로그 상시 점검"],
    ], [1.28, 2.40, 2.65, 3.05], sizes=8, row_h=0.42, grade_col=[])

    txbox(s, 0.33, 6.05, 9.3, 0.9, [
        "운영 원칙",
        "· 전역 관리자는 2~4명 이내로 제한하고 PIM 적격 할당으로 운영. 계정 잠금 대비 권한 있는 인증 관리자 최소 1명 유지",
        "· 요청 시 '수행해야 할 작업'을 먼저 확인하고, 그 작업에 필요한 최소 권한 역할을 역제안한다",
        "· 분기별 전수 조사로 미사용 할당을 회수하고, 할당 기간 만료 시 자동 회수 후 재심사한다",
    ], size=8, space=2)

    # ── 5~7. 역할 등급 분류표 ────────────────────────────────────────────
    def role_slides(grade, title_suffix, per_page=15):
        lst = [r for r in ROWS if r["grade"] == grade]
        pages = [lst[i:i + per_page] for i in range(0, len(lst), per_page)]
        for pi, page in enumerate(pages, 1):
            sfx = f" ({pi}/{len(pages)})" if len(pages) > 1 else ""
            sl = slide_frame(prs, f"[참고] 관리자 역할 등급 분류 — {grade} {title_suffix}{sfx}",
                             GRADE_DEF[grade]["desc"] + f"  ·  총 {len(lst)}건")
            rows = [["역할명 (한글)", "영문 역할명", "범주", "조합", "등급", "핵심 리스크"]]
            for r in page:
                rows.append([r["ko"], r["en"], r["category"], r["combo"],
                             r["grade"] + ("★" if r["privileged"] else ""), r["keyRisk"]])
            table(sl, 0.31, 1.34, 9.38, rows, [1.72, 2.05, 0.72, 0.55, 0.55, 3.79],
                  sizes=7, row_h=min(0.335, 5.55 / max(len(rows), 1)), grade_col=[4])
            txbox(sl, 0.33, 7.22, 9.3, 0.2,
                  "★ = Microsoft 공식 '권한 있는 역할(Privileged role)' · 조합 = 범위(S/T) + 작업(R/W) + 민감 접근(N/P)",
                  size=7, color=MUTED)

    role_slides("A3", "높음 (특권 / 전역 변경)")
    role_slides("A2", "중간 (서비스 운영 / 전역 조회)")
    role_slides("A1", "낮음 (제한적 조회)")

    p = os.path.join(DOCS, "M365 관리자 권한 분류 및 등급 기준_v1.0.pptx")
    prs.save(p)
    print("  PPT   :", os.path.basename(p), f"({len(prs.slides.__iter__.__self__._sldIdLst)}장)")


if __name__ == "__main__":
    print("M365 관리자 역할 산출물 생성")
    payload = build_json()
    inject_html(payload)
    build_excel()
    build_ppt()
    print("완료.")
